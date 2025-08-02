from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from upload_file import upload_file
from pathlib import Path
import io
import logging
from typing import List, Dict, Any

TITLE_LAYOUT = 2
SECTION_LAYOUT = 7
CONTENT_LAYOUT = 4

# Create a logger
logger = logging.getLogger(__name__)

def load_templates():
    """Loads presentation templates, returns None if not found"""
    template_4_3 = Path("template_general_4_3.pptx")
    template_16_9 = Path("template_general_16_9.pptx")
    
    # Check if templates exist
    if template_4_3.exists() and template_16_9.exists():
        return str(template_4_3), str(template_16_9)
    else:
        logger.warning("Template files not found, will use default PowerPoint templates")
        return None, None


class PowerpointPresentation:

    def __init__(self, slides: List[Dict[str, Any]], format: str):
        """Initialize PowerPoint presentation with slides and format"""

        # Validate input
        if not slides:
            raise ValueError("At least one slide is required")

        # Loads templates
        self.template_regular, self.template_wide = load_templates()

        # Create presentation based on the format used
        try:
            if format == "4:3":
                if self.template_regular:
                    self.presentation = Presentation(self.template_regular)
                else:
                    self.presentation = Presentation()  # Use default template
            elif format == "16:9":
                if self.template_wide:
                    self.presentation = Presentation(self.template_wide)
                else:
                    self.presentation = Presentation()  # Use default template
            else:
                logger.warning(f"Unknown format '{format}', defaulting to 4:3")
                if self.template_regular:
                    self.presentation = Presentation(self.template_regular)
                else:
                    self.presentation = Presentation()  # Use default template
        except Exception as e:
            logger.error(f"Failed to load template: {e}")
            logger.info("Falling back to default PowerPoint template")
            self.presentation = Presentation()  # Fallback to default template

        # Remove default slide if it exists
        if len(self.presentation.slides) > 0:
            slide_to_remove = self.presentation.slides[0]
            rId = self.presentation.slides.element.remove(slide_to_remove.element)

        # Create slides
        self._create_slides(slides)

    def _create_slides(self, slides: List[Dict[str, Any]]):
        """Create all slides from the slides data"""
        for i, slide in enumerate(slides):
            try:
                slide_type = slide.get("slide_type")

                if slide_type == "content":
                    self.create_content_slide(slide)
                elif slide_type == "section":
                    self.create_section_slide(slide)
                elif slide_type == "title":
                    self.create_title_slide(slide)
                else:
                    logger.warning(f"Unknown slide type '{slide_type}' for slide {i}, skipping")

            except Exception as e:
                logger.error(f"Failed to create slide {i}: {e}")
                raise ValueError(f"Error creating slide {i}: {str(e)}")

    def create_title_slide(self, slide: Dict[str, Any]):
        """Create a title slide"""
        try:
            title_layout = self.presentation.slide_layouts[TITLE_LAYOUT]
            title_slide = self.presentation.slides.add_slide(title_layout)

            # Set title
            if len(title_slide.placeholders) > 0:
                title_slide.placeholders[0].text = slide.get("slide_title", "")

            # Set author
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = slide.get("author", "")

        except Exception as e:
            logger.error(f"Failed to create title slide: {e}")
            raise

    def create_section_slide(self, slide: Dict[str, Any]):
        """Create a section slide"""
        try:
            section_layout = self.presentation.slide_layouts[SECTION_LAYOUT]
            section_slide = self.presentation.slides.add_slide(section_layout)

            # Set title
            if len(section_slide.placeholders) > 0:
                section_slide.placeholders[0].text = slide.get("slide_title", "")

        except Exception as e:
            logger.error(f"Failed to create section slide: {e}")
            raise

    def create_content_slide(self, slide: Dict[str, Any]):
        """Create a content slide with bullet points"""
        try:
            content_layout = self.presentation.slide_layouts[CONTENT_LAYOUT]
            content_slide = self.presentation.slides.add_slide(content_layout)

            # Set title
            if len(content_slide.placeholders) > 0:
                content_slide.placeholders[0].text = slide.get("slide_title", "")

            # Add content
            slide_text = slide.get("slide_text", [])
            if slide_text and len(content_slide.placeholders) > 1:
                # Clear existing text
                content_slide.placeholders[1].text = ""

                # Add first paragraph
                first_item = slide_text[0]
                content_slide.placeholders[1].text_frame.paragraphs[0].text = first_item.get("text", "")
                content_slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                content_slide.placeholders[1].text_frame.paragraphs[0].level = max(0, int(first_item.get("indentation_level", 1)) - 1)

                # Add remaining paragraphs
                for paragraph_data in slide_text[1:]:
                    p = content_slide.placeholders[1].text_frame.add_paragraph()
                    p.text = paragraph_data.get("text", "")
                    p.alignment = PP_ALIGN.LEFT
                    p.level = max(0, int(paragraph_data.get("indentation_level", 1)) - 1)

        except Exception as e:
            logger.error(f"Failed to create content slide: {e}")
            raise

    def save(self) -> io.BytesIO:
        """Save presentation to BytesIO object"""
        try:
            file_like_object = io.BytesIO()
            self.presentation.save(file_like_object)
            file_like_object.seek(0)
            return file_like_object
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise

def create_presentation(slides: List[Dict[str, Any]], format: str = "4:3") -> str:
    """Creates new presentation."""

    try:
        # Validate input
        if not slides:
            raise ValueError("No slides provided")

        # Create presentation
        presentation = PowerpointPresentation(slides, format)

        # Save presentation
        file_object = presentation.save()

        # Upload presentation
        text = upload_file(file_object, "pptx")
        file_object.close()

        # Return presentation link
        return text

    except Exception as e:
        logger.error(f"Failed to create presentation: {e}")
        raise
